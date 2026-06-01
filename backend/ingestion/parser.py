import fitz
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)


def parse_pdf(path: str) -> tuple[str, dict]:
    """Parse PDF and return text with page mapping. Handles errors gracefully."""
    try:
        doc = fitz.open(path)
        
        # Check if encrypted
        if doc.is_encrypted:
            logger.error("PDF is password protected")
            raise ValueError("PDF is password protected. Please provide an unlocked PDF.")
        
        if doc.page_count == 0:
            raise ValueError("PDF has no pages")
        
        pages = []
        page_map = {}
        
        for page_num, page in enumerate(doc, start=1):
            try:
                # Extract text
                text = page.get_text().strip()
                
                # Try to extract tables if text is minimal
                if len(text) < 50:
                    # Check for tables
                    tables = page.find_tables()
                    if tables:
                        for table in tables.tables:
                            try:
                                # Convert table to text representation
                                table_text = "\n".join([" | ".join(row) for row in table.extract()])
                                text += f"\n[Table]\n{table_text}\n"
                            except:
                                pass
                
                if text:
                    page_map[len(pages)] = page_num
                    pages.append(f"[Page {page_num}]\n{text}")
                else:
                    logger.warning(f"Page {page_num} has no extractable text (might be scanned)")
            
            except Exception as e:
                logger.warning(f"Error extracting page {page_num}: {e}")
                continue
        
        doc.close()
        
        if not pages:
            raise ValueError("No text could be extracted from PDF. It might be scanned or image-based.")
        
        return "\n\n".join(pages), page_map
    
    except fitz.FileDataError:
        raise ValueError("Invalid or corrupted PDF file")
    except Exception as e:
        logger.error(f"PDF parsing failed: {e}")
        raise ValueError(f"Failed to parse PDF: {str(e)}")


def parse_pptx(path: str) -> tuple[str, dict]:
    """Parse PPTX and return text with slide mapping. Handles errors gracefully."""
    try:
        prs = Presentation(path)
        
        if len(prs.slides) == 0:
            raise ValueError("PowerPoint has no slides")
        
        slides = []
        slide_map = {}
        
        for i, slide in enumerate(prs.slides, start=1):
            try:
                texts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            t = shape.text_frame.text.strip()
                            if t:
                                texts.append(t)
                        
                        # Extract text from tables
                        elif shape.has_table:
                            table = shape.table
                            table_text = []
                            for row in table.rows:
                                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                                if row_text.strip():
                                    table_text.append(row_text)
                            if table_text:
                                texts.append("[Table]\n" + "\n".join(table_text))
                    
                    except Exception as e:
                        logger.warning(f"Error extracting shape from slide {i}: {e}")
                        continue
                
                # Extract notes if available
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            texts.append(f"[Notes: {notes_text}]")
                except:
                    pass
                
                if texts:
                    slide_map[len(slides)] = i
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
                else:
                    logger.warning(f"Slide {i} has no extractable text")
            
            except Exception as e:
                logger.warning(f"Error processing slide {i}: {e}")
                continue
        
        if not slides:
            raise ValueError("No text could be extracted from PowerPoint")
        
        return "\n\n".join(slides), slide_map
    
    except Exception as e:
        logger.error(f"PPTX parsing failed: {e}")
        raise ValueError(f"Failed to parse PowerPoint: {str(e)}")


def parse_file(path: str, filename: str) -> tuple[str, dict]:
    """Parse file and return text with page/slide mapping."""
    if filename.lower().endswith(".pdf"):
        return parse_pdf(path)
    elif filename.lower().endswith(".pptx"):
        return parse_pptx(path)
    else:
        raise ValueError(f"Unsupported file format: {filename}")
